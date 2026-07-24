# -*- coding: utf-8 -*-
"""
Пакет sales_assistant.ui.main_window.

Главное окно приложения.
"""
import uuid
import os
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QDockWidget, QStackedWidget, QWidget, QVBoxLayout,
    QLabel, QMessageBox, QPushButton, QLineEdit, QFormLayout, QDialog,
    QDialogButtonBox, QTextEdit, QSizePolicy, QHBoxLayout, QSlider
)
from PySide6.QtCore import Qt, QUrl, QTimer
from PySide6.QtGui import QDesktopServices, QPixmap

from sales_assistant.db.engine import SessionLocal
from sales_assistant.services.user_service import UserService
from sales_assistant.services.company_service import CompanyService
from sales_assistant.ui.panels.navigation_panel import NavigationPanel
from sales_assistant.ui.panels.company_card import CompanyCardWidget


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Sales Assistant")
        self.resize(1200, 800)

        self.current_user = UserService.ensure_local_admin()
        self.current_user_id = self.current_user.id

        self._setup_ui()

    def _setup_ui(self):
        # Центральный стек
        self.central_stack = QStackedWidget()
        welcome_widget = QLabel("Выберите компанию из списка или добавьте новую")
        welcome_widget.setAlignment(Qt.AlignCenter)
        self.central_stack.addWidget(welcome_widget)

        self.company_card = CompanyCardWidget(self.current_user_id)
        self.company_card.file_preview_requested.connect(self._show_file_preview)
        self.central_stack.addWidget(self.company_card)
        self.setCentralWidget(self.central_stack)

        # Левая панель навигации
        nav_panel = NavigationPanel(self.current_user_id)
        nav_panel.company_selected.connect(self._show_company)
        nav_panel.add_company_requested.connect(self._add_company_dialog)
        nav_dock = QDockWidget("Навигация", self)
        nav_dock.setObjectName("NavigationDock")
        nav_dock.setWidget(nav_panel)
        nav_dock.setAllowedAreas(Qt.LeftDockWidgetArea | Qt.RightDockWidgetArea)
        self.addDockWidget(Qt.LeftDockWidgetArea, nav_dock)

        # Правая панель предпросмотра
        self.preview_dock = QDockWidget("Просмотр файла", self)
        self.preview_dock.setObjectName("PreviewDock")
        self.preview_dock.setAllowedAreas(Qt.RightDockWidgetArea | Qt.LeftDockWidgetArea)
        self.preview_dock.setMinimumWidth(200)
        self.preview_dock.setBaseSize(350, 400)

        self.preview_stack = QStackedWidget()

        # 0: приглашение
        empty_widget = QWidget()
        empty_layout = QVBoxLayout(empty_widget)
        empty_layout.addWidget(QLabel("Выберите файл для быстрого просмотра"))
        empty_layout.setAlignment(Qt.AlignCenter)
        self.preview_stack.addWidget(empty_widget)

        # 1: текст (в том числе HTML)
        self.text_preview = QTextEdit()
        self.text_preview.setReadOnly(True)
        self.text_preview.setPlaceholderText("Текстовое содержимое...")
        self.preview_stack.addWidget(self.text_preview)

        # 2: изображения / PDF
        self.image_preview = QLabel()
        self.image_preview.setAlignment(Qt.AlignCenter)
        self.image_preview.setSizePolicy(QSizePolicy.Ignored, QSizePolicy.Ignored)
        self.image_preview.setScaledContents(False)
        self.preview_stack.addWidget(self.image_preview)

        # 3: неподдерживаемый формат
        unsupported_widget = QWidget()
        unsupported_layout = QVBoxLayout(unsupported_widget)
        self.unsupported_label = QLabel("Формат не поддерживается для предпросмотра.")
        self.unsupported_label.setAlignment(Qt.AlignCenter)
        unsupported_layout.addWidget(self.unsupported_label)
        self.preview_stack.addWidget(unsupported_widget)

        # Элементы управления страницами и масштабом
        controls_widget = QWidget()
        controls_layout = QHBoxLayout(controls_widget)

        self.btn_prev_page = QPushButton("<")
        self.btn_prev_page.setToolTip("Предыдущая страница")
        self.btn_prev_page.clicked.connect(self._prev_page)
        self.btn_next_page = QPushButton(">")
        self.btn_next_page.setToolTip("Следующая страница")
        self.btn_next_page.clicked.connect(self._next_page)
        self.lbl_page = QLabel("1 / 1")
        self.lbl_page.setAlignment(Qt.AlignCenter)

        self.zoom_slider = QSlider(Qt.Horizontal)
        self.zoom_slider.setRange(50, 300)
        self.zoom_slider.setValue(100)
        self.zoom_slider.setTickInterval(50)
        self.zoom_slider.setTickPosition(QSlider.TicksBelow)
        self.zoom_slider.valueChanged.connect(self._zoom_changed)
        self.zoom_label = QLabel("100%")

        controls_layout.addWidget(self.btn_prev_page)
        controls_layout.addWidget(self.lbl_page)
        controls_layout.addWidget(self.btn_next_page)
        controls_layout.addWidget(self.zoom_label)
        controls_layout.addWidget(self.zoom_slider)

        # Кнопка "Открыть в системе"
        open_btn = QPushButton("Открыть в системе")
        open_btn.clicked.connect(self._open_current_file_external)

        container = QWidget()
        container_layout = QVBoxLayout(container)
        container_layout.addWidget(self.preview_stack)
        container_layout.addWidget(controls_widget)
        container_layout.addWidget(open_btn)

        self.preview_dock.setWidget(container)
        self.preview_dock.hide()
        self.addDockWidget(Qt.RightDockWidgetArea, self.preview_dock)

        # Внутренние переменные
        self.current_preview_file = None
        self._pdf_doc = None
        self._pdf_page_count = 0
        self._current_page_idx = 0
        self._zoom_factor = 1.0
        self._base_pixmap = None

        # Скрыть элементы управления PDF по умолчанию
        self._hide_pdf_controls()

    # --------------------------------------------------------------
    # Предпросмотр файлов
    # --------------------------------------------------------------
    def _show_file_preview(self, file_path: str):
        """Загружает файл в панель предпросмотра."""
        self.current_preview_file = file_path
        if not file_path:
            self.preview_stack.setCurrentIndex(0)
            self.preview_dock.hide()
            return

        if not os.path.exists(file_path):
            self.preview_stack.setCurrentIndex(0)
            self.preview_dock.show()
            return

        ext = os.path.splitext(file_path)[1].lower()
        text_exts = {'.txt', '.py', '.md', '.csv', '.log', '.json', '.xml', '.html', '.css', '.js', '.ini', '.cfg', '.yaml', '.yml'}
        image_exts = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.svg', '.ico', '.webp'}
        pdf_exts = {'.pdf'}
        office_exts = {'.docx', '.xlsx', '.pptx', '.odt', '.ods', '.odp'}

        if self._pdf_doc:
            self._pdf_doc.close()
            self._pdf_doc = None

        if ext in text_exts:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                self.text_preview.setPlainText(content)
                self.preview_stack.setCurrentIndex(1)
                self._hide_pdf_controls()
            except Exception as e:
                self._show_unsupported(str(e))
                self._hide_pdf_controls()

        elif ext in image_exts:
            self._load_image(file_path)
            self._hide_pdf_controls()

        elif ext in pdf_exts:
            self._load_pdf(file_path)
            self._show_pdf_controls()

        elif ext in office_exts:
            self._load_office_text(file_path)
            self._hide_pdf_controls()

        else:
            self._show_unsupported("Формат не поддерживается.")
            self._hide_pdf_controls()

        self.preview_dock.show()
        QApplication.processEvents()
        self._render_image_at_zoom()

    def _load_image(self, file_path: str):
        pixmap = QPixmap(file_path)
        if pixmap.isNull():
            self._show_unsupported("Не удалось загрузить изображение.")
            return
        self._base_pixmap = pixmap
        self.preview_stack.setCurrentIndex(2)

    def _load_pdf(self, file_path: str):
        try:
            import fitz
        except ImportError:
            self._show_unsupported("Для предпросмотра PDF установите PyMuPDF\n(pip install PyMuPDF)")
            return
        try:
            doc = fitz.open(file_path)
            self._pdf_doc = doc
            self._pdf_page_count = doc.page_count
            self._current_page_idx = 0
            self._render_pdf_page()
            self.preview_stack.setCurrentIndex(2)
            self._update_page_label()
        except Exception as e:
            self._show_unsupported(f"Ошибка открытия PDF:\n{e}")

    def _load_office_text(self, file_path: str):
        """Извлекает форматированный текст из офисных документов и показывает как HTML."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == '.docx':
                html = self._docx_to_html(file_path)
                self.text_preview.setHtml(html)
                self.preview_stack.setCurrentIndex(1)
            elif ext == '.xlsx':
                html = self._xlsx_to_html(file_path)
                self.text_preview.setHtml(html)
                self.preview_stack.setCurrentIndex(1)
            elif ext == '.pptx':
                html = self._pptx_to_html(file_path)
                self.text_preview.setHtml(html)
                self.preview_stack.setCurrentIndex(1)
            elif ext in ('.odt', '.ods', '.odp'):
                text = self._odf_to_text(file_path, ext)
                if not text:
                    self._show_unsupported("Документ не содержит текста или формат не распознан.")
                else:
                    self.text_preview.setPlainText(text)
                    self.preview_stack.setCurrentIndex(1)
            else:
                self._show_unsupported("Неподдерживаемый формат.")
        except ImportError as e:
            self._show_unsupported(f"Не установлена библиотека для {ext}.\n" + "pip install python-docx openpyxl python-pptx odfpy")
        except Exception as e:
            self._show_unsupported(f"Ошибка чтения документа:\n{e}")

    def _docx_to_html(self, file_path: str) -> str:
        """Преобразует DOCX в HTML с базовым форматированием."""
        import docx
        doc = docx.Document(file_path)
        html_parts = ["<html><body style='font-size:12pt; font-family:Segoe UI;'>"]

        for para in doc.paragraphs:
            # Определяем стиль параграфа
            if para.style.name.startswith('Heading'):
                level = para.style.name.split()[-1]
                tag = f"h{level}"
                html_parts.append(f"<{tag}>")
            else:
                html_parts.append("<p>")

            # Проходим по runs внутри параграфа
            for run in para.runs:
                text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                if run.bold:
                    text = f"<b>{text}</b>"
                if run.italic:
                    text = f"<i>{text}</i>"
                if run.underline:
                    text = f"<u>{text}</u>"
                html_parts.append(text)
            html_parts.append(f"</{tag if para.style.name.startswith('Heading') else 'p'}>")

        html_parts.append("</body></html>")
        return '\n'.join(html_parts)

    def _xlsx_to_html(self, file_path: str) -> str:
        """Преобразует XLSX в HTML-таблицу."""
        import openpyxl
        from openpyxl.utils import get_column_letter

        wb = openpyxl.load_workbook(file_path, data_only=True)
        html_parts = ["<html><body style='font-size:11pt; font-family:Segoe UI;'>"]

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            html_parts.append(f"<h3>{sheet_name}</h3>")
            html_parts.append("<table border='1' cellpadding='4' cellspacing='0' style='border-collapse: collapse; width: 100%;'>")

            for row in ws.iter_rows():
                html_parts.append("<tr>")
                for cell in row:
                    value = str(cell.value) if cell.value is not None else ''
                    # Можно добавить жирность для первой строки, если захотим
                    if cell.row == 1:
                        html_parts.append(f"<td style='font-weight: bold; background-color: #f0f0f0;'>{value}</td>")
                    else:
                        html_parts.append(f"<td>{value}</td>")
                html_parts.append("</tr>")
            html_parts.append("</table><br>")
        html_parts.append("</body></html>")
        return '\n'.join(html_parts)

    def _pptx_to_html(self, file_path: str) -> str:
        """Преобразует PPTX в HTML, группируя текст по слайдам."""
        import pptx
        prs = pptx.Presentation(file_path)
        html_parts = ["<html><body style='font-size:12pt; font-family:Segoe UI;'>"]

        for slide_num, slide in enumerate(prs.slides, start=1):
            html_parts.append(f"<h2>Слайд {slide_num}</h2>")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        runs_html = []
                        for run in para.runs:
                            text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            if run.font.bold:
                                text = f"<b>{text}</b>"
                            if run.font.italic:
                                text = f"<i>{text}</i>"
                            runs_html.append(text)
                        para_text = ''.join(runs_html)
                        if para_text.strip():
                            html_parts.append(f"<p>{para_text}</p>")
            html_parts.append("<hr>")
        html_parts.append("</body></html>")
        return '\n'.join(html_parts)

    def _odf_to_text(self, file_path: str, ext: str) -> str:
        """Извлекает plain text из ODF-документов."""
        import odf.opendocument as od
        if ext == '.odt':
            from odf import text as odf_text
            doc = od.load(file_path)
            content = doc.getElementsByType(odf_text.P)
            return '\n'.join([str(c) for c in content])
        elif ext == '.ods':
            from odf.table import Table, TableRow, TableCell
            doc = od.load(file_path)
            tables = doc.getElementsByType(Table)
            sheets = []
            for table in tables:
                rows = []
                for row in table.getElementsByType(TableRow):
                    cells = [str(cell) for cell in row.getElementsByType(TableCell)]
                    rows.append('\t'.join(cells))
                sheets.append('\n'.join(rows))
            return '\n\n'.join(sheets)
        elif ext == '.odp':
            from odf.draw import Page
            doc = od.load(file_path)
            pages = doc.getElementsByType(Page)
            slides = [str(page) for page in pages]
            return '\n\n'.join(slides)
        return ""

    def _show_unsupported(self, message: str):
        self.unsupported_label.setText(message)
        self.preview_stack.setCurrentIndex(3)

    def _render_pdf_page(self):
        if not self._pdf_doc:
            return
        page = self._pdf_doc[self._current_page_idx]
        dpi = 150 * self._zoom_factor
        pix = page.get_pixmap(dpi=int(dpi))
        img_data = pix.tobytes("ppm")
        pixmap = QPixmap()
        pixmap.loadFromData(img_data, "PPM")
        self._base_pixmap = pixmap
        self._render_image_at_zoom()

    def _render_image_at_zoom(self):
        if self._base_pixmap is None or self._base_pixmap.isNull():
            return
        if self.image_preview.width() > 0 and self.image_preview.height() > 0:
            target_w = self.image_preview.width()
            target_h = self.image_preview.height()
        else:
            target_w = self._base_pixmap.width() * self._zoom_factor / 100.0
            target_h = self._base_pixmap.height() * self._zoom_factor / 100.0

        scaled = self._base_pixmap.scaled(
            int(target_w), int(target_h),
            Qt.KeepAspectRatio, Qt.SmoothTransformation
        )
        self.image_preview.setPixmap(scaled)

    def _prev_page(self):
        if self._pdf_doc and self._current_page_idx > 0:
            self._current_page_idx -= 1
            self._render_pdf_page()
            self._update_page_label()

    def _next_page(self):
        if self._pdf_doc and self._current_page_idx < self._pdf_page_count - 1:
            self._current_page_idx += 1
            self._render_pdf_page()
            self._update_page_label()

    def _update_page_label(self):
        self.lbl_page.setText(f"{self._current_page_idx + 1} / {self._pdf_page_count}")

    def _zoom_changed(self, value):
        self._zoom_factor = value / 100.0
        self.zoom_label.setText(f"{value}%")
        if self._pdf_doc:
            self._render_pdf_page()
        elif self._base_pixmap:
            self._render_image_at_zoom()

    def _hide_pdf_controls(self):
        self.btn_prev_page.hide()
        self.btn_next_page.hide()
        self.lbl_page.hide()

    def _show_pdf_controls(self):
        self.btn_prev_page.show()
        self.btn_next_page.show()
        self.lbl_page.show()

    def _open_current_file_external(self):
        if self.current_preview_file and os.path.exists(self.current_preview_file):
            try:
                os.startfile(self.current_preview_file)
            except Exception:
                QDesktopServices.openUrl(QUrl.fromLocalFile(self.current_preview_file))

    def resizeEvent(self, event):
        super().resizeEvent(event)
        if self.preview_stack.currentIndex() == 2 and self._base_pixmap:
            self._render_image_at_zoom()

    # --------------------------------------------------------------
    # Навигация и компании
    # --------------------------------------------------------------
    def _show_company(self, company_id: uuid.UUID):
        self.company_card.load_company(company_id)
        self.central_stack.setCurrentIndex(1)

    def _add_company_dialog(self):
        dialog = QDialog(self)
        dialog.setWindowTitle("Новая компания")
        layout = QFormLayout(dialog)
        name_edit = QLineEdit()
        layout.addRow("Название:", name_edit)
        buttons = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        buttons.accepted.connect(dialog.accept)
        buttons.rejected.connect(dialog.reject)
        layout.addWidget(buttons)

        if dialog.exec() == QDialog.Accepted:
            name = name_edit.text().strip()
            if name:
                with SessionLocal() as session:
                    try:
                        service = CompanyService(self.current_user_id)
                        company = service.create(session, name)
                        session.commit()
                        nav_dock = self.findChild(QDockWidget, "NavigationDock")
                        if nav_dock:
                            nav_panel = nav_dock.widget()
                            nav_panel._load_companies()
                        self._show_company(company.id)
                    except Exception as e:
                        QMessageBox.critical(self, "Ошибка", f"Не удалось создать компанию:\n{e}")
                    finally:
                        session.close()
            else:
                QMessageBox.warning(self, "Предупреждение", "Название компании не может быть пустым.")